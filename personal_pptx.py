from personal_pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Define slide layout (Title + Content)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# --- Slide 1: Title ---
slide1 = prs.slides.add_slide(title_slide_layout)
slide1.shapes.title.text = "Personal Pitch — Humphrey Maina"
slide1.placeholders[1].text = ""  # no subtitle

# --- Slide 2: About Me ---
slide2 = prs.slides.add_slide(content_slide_layout)
slide2.shapes.title.text = "About Me"
content2 = slide2.placeholders[1].text_frame
content2.text = (
    "• Passionate about technology and innovation\n"
    "• Currently pursuing AI for Software Engineering at Mount Kenya University\n"
    "• Power Learn Project Africa scholar\n"
    "• Strong interest in Data Science, Machine Learning, and real-world AI solutions"
)

# --- Slide 3: Skills & Projects ---
slide3 = prs.slides.add_slide(content_slide_layout)
slide3.shapes.title.text = "Skills & Projects"
content3 = slide3.placeholders[1].text_frame
content3.text = (
    "• Programming: Python, SQL, TensorFlow, Pandas\n"
    "• Data Analysis & Visualization: Power BI, Matplotlib, Excel\n"
    "• AI & ML: Model development, data preprocessing, automation\n"
    "• Coursework and hands-on projects in AI applications"
)

# --- Slide 4: Vision & Goals ---
slide4 = prs.slides.add_slide(content_slide_layout)
slide4.shapes.title.text = "Vision & Goals"
content4 = slide4.placeholders[1].text_frame
content4.text = (
    "• Aspire to become a leading Data Scientist in Africa\n"
    "• Apply AI to solve challenges in education, healthcare, and sustainability\n"
    "• Commit to continuous learning, collaboration, and innovation"
)

# Save presentation
prs.save("Personal_Pitch_Humphrey_Maina.pptx")
print("✅ Presentation saved as Personal_Pitch_Humphrey_Maina.pptx")